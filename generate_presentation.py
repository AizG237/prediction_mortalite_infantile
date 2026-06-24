#!/usr/bin/env python3
"""
Génération automatique de la présentation PowerPoint
Analyse des Facteurs Associés à la Mortalité Infantile au Cameroun
Auteur   : M. Yannick FONO BINDJEME
Directeur: Pr Nguefack-Tsague Georges
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ─────────────────────────────────────────────────────────────────────────────
# PALETTE (identique au thème Streamlit)
# ─────────────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x3A, 0x5C)
BLUE    = RGBColor(0x2C, 0x5F, 0x8A)
SKY     = RGBColor(0x34, 0x98, 0xDB)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF8, 0xF9, 0xFA)
GREEN   = RGBColor(0x27, 0xAE, 0x60)
ORANGE  = RGBColor(0xF3, 0x9C, 0x12)
RED     = RGBColor(0xE7, 0x4C, 0x3C)
GOLD    = RGBColor(0xE8, 0xC0, 0x45)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
GRAY    = RGBColor(0x6C, 0x75, 0x7D)
LGRAY   = RGBColor(0xDE, 0xE2, 0xE6)
CARD_BG = RGBColor(0xEA, 0xF2, 0xFB)
CARD_G  = RGBColor(0xE6, 0xF9, 0xF0)
CARD_O  = RGBColor(0xFE, 0xF6, 0xE4)
CARD_R  = RGBColor(0xFD, 0xED, 0xEB)
NAVY2   = RGBColor(0x0D, 0x1B, 0x2A)

# ─────────────────────────────────────────────────────────────────────────────
# DIMENSIONS
# ─────────────────────────────────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)
N  = 19   # nombre total de diapositives

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

_pg = [0]

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

def img_path(rel):
    return os.path.join(BASE_DIR, rel)

# ─────────────────────────────────────────────────────────────────────────────
# UTILITAIRES
# ─────────────────────────────────────────────────────────────────────────────

def new_slide(bg=BG):
    _pg[0] += 1
    sl = prs.slides.add_slide(BLANK)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = bg
    return sl, _pg[0]


def rect(sl, l, t, w, h, fill, line=None, lw=1.0, rnd=False):
    shp = sl.shapes.add_shape(5 if rnd else 1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp


def txt(sl, text, l, t, w, h, sz=12, bold=False, italic=False,
        c=DARK, align=PP_ALIGN.LEFT, wrap=True):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = c
    return bx


def para(tf, text, sz=11, bold=False, italic=False, c=DARK,
         align=PP_ALIGN.LEFT, bullet=True, symbol='▸'):
    p  = tf.add_paragraph()
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size  = Pt(sz)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = c
    if bullet:
        pPr = p._p.get_or_add_pPr()
        bu  = etree.SubElement(pPr, qn('a:buChar'))
        bu.set('char', symbol)
    return p


def img(sl, rel_path, l, t, w, h=None):
    p = img_path(rel_path)
    if os.path.exists(p):
        if h:
            return sl.shapes.add_picture(p, l, t, w, h)
        return sl.shapes.add_picture(p, l, t, w)
    return None


def foot(sl, n):
    rect(sl, 0, Inches(7.2), SW, Inches(0.3), NAVY)
    txt(sl,
        "M. Yannick FONO BINDJEME  |  Institut Saint Jean Ingénieur  |  EDS Cameroun 2018",
        Inches(0.3), Inches(7.22), Inches(11.5), Inches(0.26),
        sz=7.5, c=WHITE)
    txt(sl, f"{n} / {N}",
        Inches(11.9), Inches(7.22), Inches(1.2), Inches(0.26),
        sz=8, bold=True, c=GOLD, align=PP_ALIGN.RIGHT)


def head(sl, title, sub=None, badge=None, badge_c=BLUE):
    rect(sl, 0, 0, SW, Inches(1.3), NAVY)
    rect(sl, 0, Inches(1.3), SW, Inches(0.05), GOLD)
    txt(sl, title, Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.65),
        sz=22, bold=True, c=WHITE)
    if sub:
        txt(sl, sub, Inches(0.5), Inches(0.8), Inches(11), Inches(0.42),
            sz=12, italic=True, c=RGBColor(0xAE, 0xC6, 0xE0))
    if badge:
        rect(sl, Inches(11.83), Inches(0.22), Inches(1.28), Inches(0.3),
             badge_c, rnd=True)
        txt(sl, badge,
            Inches(11.83), Inches(0.24), Inches(1.28), Inches(0.3),
            sz=8.5, bold=True, c=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# ANIMATIONS FADE-IN (cascade automatique au chargement de la diapositive)
# ─────────────────────────────────────────────────────────────────────────────

def auto_fade(sl, shapes, start_delay=300, step=220, duration=550):
    """Chaque forme apparaît en fondu après un délai croissant."""
    if not shapes:
        return
    NS  = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NSA = "http://schemas.openxmlformats.org/drawingml/2006/main"

    root_xml = (
        f'<p:timing xmlns:p="{NS}" xmlns:a="{NSA}">'
        '<p:tnLst><p:par>'
        '<p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">'
        '<p:childTnLst/>'
        '</p:cTn></p:par></p:tnLst>'
        '<p:bldLst/></p:timing>'
    )
    root = etree.fromstring(root_xml)
    cl   = root.find(f".//{{{NS}}}childTnLst")

    aid = 2
    for i, shp in enumerate(shapes):
        sid   = str(shp.shape_id)
        delay = start_delay + i * step
        dur   = str(duration)
        blk = (
            f'<p:par xmlns:p="{NS}" xmlns:a="{NSA}">'
            f'<p:cTn id="{aid}" fill="hold">'
            f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
            f'<p:childTnLst><p:par>'
            f'<p:cTn id="{aid+1}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>'
            f'<p:set><p:cBhvr>'
            f'<p:cTn id="{aid+2}" dur="1" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
            f'<p:animEffect transition="in" filter="fade">'
            f'<p:cBhvr><p:cTn id="{aid+3}" dur="{dur}"/>'
            f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
            f'</p:cBhvr></p:animEffect>'
            f'</p:childTnLst></p:cTn>'
            f'</p:par></p:childTnLst></p:cTn></p:par>'
        )
        cl.append(etree.fromstring(blk))
        aid += 4

    sl._element.append(root)


def click_fade(sl, shapes, duration=550):
    """Chaque forme apparaît au clic suivant."""
    if not shapes:
        return
    NS  = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NSA = "http://schemas.openxmlformats.org/drawingml/2006/main"

    root_xml = (
        f'<p:timing xmlns:p="{NS}" xmlns:a="{NSA}">'
        '<p:tnLst><p:par>'
        '<p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">'
        '<p:childTnLst>'
        '<p:seq concurrent="1" nextAc="seek">'
        '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst/></p:cTn>'
        '<p:prevCondLst><p:cond evt="onPrevClick" delay="0"><p:tn/></p:cond></p:prevCondLst>'
        '<p:nextCondLst><p:cond evt="onNextClick" delay="0"><p:tn/></p:cond></p:nextCondLst>'
        '</p:seq>'
        '</p:childTnLst></p:cTn></p:par></p:tnLst>'
        '<p:bldLst/></p:timing>'
    )
    root = etree.fromstring(root_xml)
    cln  = root.find(f".//{{{NS}}}cTn[@nodeType='mainSeq']/{{{NS}}}childTnLst")

    aid = 3
    for i, shp in enumerate(shapes):
        sid  = str(shp.shape_id)
        dl   = "indefinite" if i == 0 else "200"
        dur_ = str(duration)
        blk  = (
            f'<p:par xmlns:p="{NS}" xmlns:a="{NSA}">'
            f'<p:cTn id="{aid}" fill="hold">'
            f'<p:stCondLst><p:cond delay="{dl}"/></p:stCondLst>'
            f'<p:childTnLst><p:par>'
            f'<p:cTn id="{aid+1}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>'
            f'<p:set><p:cBhvr>'
            f'<p:cTn id="{aid+2}" dur="1" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
            f'<p:animEffect transition="in" filter="fade">'
            f'<p:cBhvr><p:cTn id="{aid+3}" dur="{dur_}"/>'
            f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
            f'</p:cBhvr></p:animEffect>'
            f'</p:childTnLst></p:cTn>'
            f'</p:par></p:childTnLst></p:cTn></p:par>'
        )
        cln.append(etree.fromstring(blk))
        aid += 4

    sl._element.append(root)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 1  –  PAGE DE TITRE
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide(bg=NAVY2)

# Fond dégradé simulé par deux rectangles
rect(sl, 0, 0, SW, SH, NAVY)
rect(sl, 0, 0, SW, Inches(3.2), NAVY2)

# Bande accent dorée
rect(sl, 0, Inches(3.1), SW, Inches(0.06), GOLD)

# Logo/icône symbolique (carré bleu décoratif haut-gauche)
rect(sl, Inches(0.4), Inches(0.4), Inches(0.12), Inches(2.5), GOLD)
rect(sl, Inches(0.62), Inches(0.4), Inches(0.06), Inches(2.5),
     RGBColor(0x4A, 0x8F, 0xC4))

# Titre principal
b1 = txt(sl, "Analyse des Facteurs Associés à la",
         Inches(1.0), Inches(0.4), Inches(11.8), Inches(0.8),
         sz=30, bold=True, c=WHITE)

b2 = txt(sl, "Mortalité Infantile au Cameroun",
         Inches(1.0), Inches(1.1), Inches(11.8), Inches(0.85),
         sz=34, bold=True, c=GOLD)

b3 = txt(sl,
         "Approche par Régression Logistique Pondérée et Machine Learning",
         Inches(1.0), Inches(1.95), Inches(11.8), Inches(0.5),
         sz=14, italic=True, c=RGBColor(0xAE, 0xC6, 0xE0))

# Séparateur
rect(sl, Inches(1.0), Inches(3.22), Inches(11.3), Inches(0.03),
     RGBColor(0x4A, 0x8F, 0xC4))

# Carte auteur
rect(sl, Inches(1.0), Inches(3.42), Inches(5.4), Inches(3.4),
     RGBColor(0x1F, 0x4B, 0x78), rnd=True)

txt(sl, "Présenté par",
    Inches(1.25), Inches(3.6), Inches(5.0), Inches(0.35),
    sz=10, italic=True, c=RGBColor(0xAE, 0xC6, 0xE0))

txt(sl, "M. Yannick FONO BINDJEME",
    Inches(1.25), Inches(3.92), Inches(5.0), Inches(0.55),
    sz=16, bold=True, c=GOLD)

txt(sl, "Étudiant en Statistiques et Informatique",
    Inches(1.25), Inches(4.45), Inches(5.0), Inches(0.4),
    sz=10.5, c=WHITE)

txt(sl, "Institut Saint Jean Ingénieur",
    Inches(1.25), Inches(4.83), Inches(5.0), Inches(0.4),
    sz=10.5, bold=True, c=SKY)

# Carte superviseur
rect(sl, Inches(7.0), Inches(3.42), Inches(5.4), Inches(3.4),
     RGBColor(0x1F, 0x4B, 0x78), rnd=True)

txt(sl, "Sous la supervision de",
    Inches(7.25), Inches(3.6), Inches(5.0), Inches(0.35),
    sz=10, italic=True, c=RGBColor(0xAE, 0xC6, 0xE0))

txt(sl, "Pr Nguefack-Tsague Georges",
    Inches(7.25), Inches(3.92), Inches(5.0), Inches(0.55),
    sz=15.5, bold=True, c=GOLD)

txt(sl, "Professeur de Biostatistiques",
    Inches(7.25), Inches(4.45), Inches(5.0), Inches(0.4),
    sz=10.5, c=WHITE)

# Données sources
rect(sl, Inches(1.0), Inches(6.72), Inches(11.3), Inches(0.38),
     RGBColor(0x1F, 0x4B, 0x78))
txt(sl,
    "Données : Enquête Démographique et de Santé (EDS) Cameroun 2018   —   Juin 2026",
    Inches(1.15), Inches(6.76), Inches(11.0), Inches(0.34),
    sz=10, italic=True, c=RGBColor(0xAE, 0xC6, 0xE0), align=PP_ALIGN.CENTER)

auto_fade(sl, [b1, b2, b3], start_delay=200, step=300, duration=800)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 2  –  SOMMAIRE
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Plan de la Présentation",
     sub="Structure en deux grandes parties : Statistiques et Machine Learning")

# Colonne gauche – Partie I
rect(sl, Inches(0.4), Inches(1.6), Inches(6.0), Inches(5.3), CARD_BG, rnd=True)
rect(sl, Inches(0.4), Inches(1.6), Inches(6.0), Inches(0.58), BLUE, rnd=True)

txt(sl, "PARTIE I", Inches(0.6), Inches(1.63), Inches(2.5), Inches(0.35),
    sz=9, bold=True, c=GOLD)
txt(sl, "Statistiques",
    Inches(0.6), Inches(1.88), Inches(5.6), Inches(0.35),
    sz=14, bold=True, c=WHITE)

items_stat = [
    "I.1  Contexte et problématique",
    "I.2  Source de données — EDS Cameroun 2018",
    "I.3  Préparation et nettoyage des données",
    "I.4  Analyse descriptive et bivariée",
    "I.5  Régression logistique pondérée",
    "I.6  Résultats et odds ratios",
]
bx_s = sl.shapes.add_textbox(Inches(0.6), Inches(2.35), Inches(5.6), Inches(4.3))
tf_s = bx_s.text_frame
tf_s.word_wrap = True
for i, item in enumerate(items_stat):
    if i == 0:
        p = tf_s.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = item
        r.font.size = Pt(11.5); r.font.color.rgb = DARK
    else:
        para(tf_s, item, sz=11.5, c=DARK, bullet=True, symbol='●')

# Colonne droite – Partie II
rect(sl, Inches(6.9), Inches(1.6), Inches(6.0), Inches(5.3), CARD_BG, rnd=True)
rect(sl, Inches(6.9), Inches(1.6), Inches(6.0), Inches(0.58), NAVY, rnd=True)

txt(sl, "PARTIE II", Inches(7.1), Inches(1.63), Inches(2.5), Inches(0.35),
    sz=9, bold=True, c=GOLD)
txt(sl, "Machine Learning",
    Inches(7.1), Inches(1.88), Inches(5.6), Inches(0.35),
    sz=14, bold=True, c=WHITE)

items_ml = [
    "II.1  Pipeline de prétraitement",
    "II.2  Six algorithmes comparés",
    "II.3  Courbes d'apprentissage",
    "II.4  Comparaison des performances",
    "II.5  Modèle champion — XGBoost & SHAP",
    "II.6  Application Streamlit et conclusion",
]
bx_m = sl.shapes.add_textbox(Inches(7.1), Inches(2.35), Inches(5.6), Inches(4.3))
tf_m = bx_m.text_frame
tf_m.word_wrap = True
for i, item in enumerate(items_ml):
    if i == 0:
        p = tf_m.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = item
        r.font.size = Pt(11.5); r.font.color.rgb = DARK
    else:
        para(tf_m, item, sz=11.5, c=DARK, bullet=True, symbol='●')

foot(sl, pn)
click_fade(sl, [bx_s, bx_m])


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 3  –  CONTEXTE ET PROBLÉMATIQUE
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Contexte et Problématique",
     sub="La mortalité infantile, un défi persistant en Afrique subsaharienne")

# Trois cartes de contexte
cards = [
    ("Contexte mondial",
     GREEN,
     CARD_G,
     ["Objectif de Développement Durable 3.2 :",
      "réduire la mortalité infantile à moins de 25 pour 1 000 naissances vivantes d'ici 2030.",
      "L'Afrique subsaharienne concentre plus de 50 % des décès d'enfants de moins de 5 ans."]),
    ("Situation au Cameroun",
     ORANGE,
     CARD_O,
     ["Taux de mortalité infantile : environ 50 pour 1 000 NV (EDS 2018).",
      "Forte disparité régionale et socio-économique.",
      "Accès limité aux soins dans les zones rurales et péri-urbaines."]),
    ("Problématique et objectif",
     BLUE,
     CARD_BG,
     ["Quels sont les facteurs déterminants de la mortalité infantile au Cameroun ?",
      "Objectif : identifier et quantifier ces facteurs par approche mixte.",
      "Méthodes : régression logistique pondérée + algorithmes de Machine Learning."]),
]

anims = []
for i, (title, accent, bg_c, bullets) in enumerate(cards):
    cx = Inches(0.35 + i * 4.32)
    bx = rect(sl, cx, Inches(1.55), Inches(4.1), Inches(5.3), bg_c, rnd=True)
    rect(sl, cx, Inches(1.55), Inches(4.1), Inches(0.06), accent)
    anims.append(bx)
    txt(sl, title, cx + Inches(0.18), Inches(1.7), Inches(3.8), Inches(0.4),
        sz=12, bold=True, c=accent)
    bx2 = sl.shapes.add_textbox(cx + Inches(0.18), Inches(2.2),
                                 Inches(3.8), Inches(4.5))
    tf = bx2.text_frame; tf.word_wrap = True
    for j, b in enumerate(bullets):
        if j == 0:
            p = tf.paragraphs[0]; r = p.add_run(); r.text = b
            r.font.size = Pt(10.5); r.font.color.rgb = DARK
        else:
            para(tf, b, sz=10.5, c=DARK, bullet=True, symbol='▸')
    anims.append(bx2)

foot(sl, pn)
auto_fade(sl, anims, start_delay=300, step=180)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 4  –  SOURCE DE DONNÉES
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Source de Données",
     sub="Enquête Démographique et de Santé — Cameroun 2018")

# Bloc gauche descriptif
rect(sl, Inches(0.4), Inches(1.6), Inches(5.5), Inches(5.3), CARD_BG, rnd=True)
rect(sl, Inches(0.4), Inches(1.6), Inches(5.5), Inches(0.55), BLUE, rnd=True)
txt(sl, "Présentation de l'EDS Cameroun 2018",
    Inches(0.6), Inches(1.65), Inches(5.0), Inches(0.45),
    sz=11.5, bold=True, c=WHITE)

b_desc = sl.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(5.1), Inches(4.4))
tf_d = b_desc.text_frame; tf_d.word_wrap = True
desc_items = [
    "Enquête nationale représentative (EDS Phase 7)",
    "Population : femmes de 15 à 49 ans",
    "Fichier utilisé : Dossier individuel femmes (CMIR71FL.dta)",
    "Sondage stratifié à deux degrés (grappes et ménages)",
    "Poids d'enquête intégrés (variable V005 / 10⁶)",
    "Plus de 1 000 variables collectées",
    "Données disponibles auprès du programme DHS",
]
for j, item in enumerate(desc_items):
    if j == 0:
        p = tf_d.paragraphs[0]; r = p.add_run(); r.text = item
        r.font.size = Pt(11); r.font.color.rgb = DARK
    else:
        para(tf_d, item, sz=11, c=DARK, bullet=True, symbol='▸')

# Métriques à droite
metrics = [
    ("14 677", "Femmes dans l'analyse finale", BLUE, WHITE),
    ("10 494", "Femmes avec au moins un enfant né vivant\n(sous-échantillon de régression)", NAVY, WHITE),
    ("~34 %", "Prévalence de la mortalité infantile\n(au moins un enfant décédé)", RED, WHITE),
    ("20 +", "Variables prédictives retenues\n(sociodémographiques, reproductives, soins)", GREEN, WHITE),
]
for i, (val, lab, c1, c2) in enumerate(metrics):
    cy = Inches(1.6 + i * 1.32)
    bk = rect(sl, Inches(6.3), cy, Inches(6.6), Inches(1.18), c1, rnd=True)
    txt(sl, val, Inches(6.55), cy + Inches(0.12), Inches(2.0), Inches(0.58),
        sz=28, bold=True, c=GOLD)
    txt(sl, lab, Inches(8.65), cy + Inches(0.18), Inches(4.0), Inches(0.82),
        sz=10.5, c=c2, wrap=True)

foot(sl, pn)
click_fade(sl, [b_desc])


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 5  –  SÉPARATEUR PARTIE I
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide(bg=NAVY)

# Rectangles décoratifs
rect(sl, 0, 0, SW, SH, NAVY)
rect(sl, 0, Inches(3.1), SW, Inches(1.3), BLUE)
rect(sl, 0, Inches(3.1), SW, Inches(0.06), GOLD)
rect(sl, 0, Inches(4.34), SW, Inches(0.06), GOLD)

# Numéro de partie – grand
txt(sl, "I", Inches(0.5), Inches(0.5), Inches(2.5), Inches(3.5),
    sz=180, bold=True, c=RGBColor(0x2C, 0x5F, 0x8A), align=PP_ALIGN.LEFT)

b_part = txt(sl, "PARTIE I",
             Inches(3.0), Inches(3.15), Inches(9.0), Inches(0.55),
             sz=22, bold=True, c=GOLD)

b_titre = txt(sl, "Analyse Statistique",
              Inches(3.0), Inches(3.62), Inches(9.0), Inches(0.72),
              sz=38, bold=True, c=WHITE)

txt(sl,
    "Exploration des données  •  Analyse bivariée  •  Régression logistique pondérée  •  Résultats",
    Inches(3.0), Inches(4.45), Inches(9.8), Inches(0.42),
    sz=12, italic=True, c=RGBColor(0xAE, 0xC6, 0xE0))

foot(sl, pn)
auto_fade(sl, [b_part, b_titre], start_delay=200, step=400)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 6  –  PRÉPARATION DES DONNÉES
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Préparation et Nettoyage des Données",
     sub="Construction de la variable cible et traitement des données manquantes",
     badge="PARTIE I", badge_c=BLUE)

# Pipeline en flèches
steps = [
    ("1. Chargement", BLUE,
     "Fichier Stata CMIR71FL.dta\nSélection de 22 variables\n1 000+ variables disponibles"),
    ("2. Nettoyage", NAVY,
     "Codes EDS (9, 99, 998, 999)\ntraités comme valeurs manquantes\nImputation médiane / mode"),
    ("3. Cible Y", RGBColor(0x7B, 0x2D, 0x8B),
     "Y = 1 si V206 + V207 > 0\n(au moins un enfant décédé)\nNulligestes → Y = 0"),
    ("4. Datasets", GREEN,
     "df clean : toutes les femmes\ndf stat : 10 494 mères\ndf ml  : 14 677 femmes"),
]
for i, (title, c, desc) in enumerate(steps):
    cx = Inches(0.35 + i * 3.22)
    bk = rect(sl, cx, Inches(1.65), Inches(3.0), Inches(2.05), c, rnd=True)
    txt(sl, title, cx + Inches(0.15), Inches(1.75), Inches(2.75), Inches(0.42),
        sz=12, bold=True, c=GOLD)
    bx2 = sl.shapes.add_textbox(cx + Inches(0.15), Inches(2.2), Inches(2.75), Inches(1.4))
    tf2 = bx2.text_frame; tf2.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            p = tf2.paragraphs[0]; r = p.add_run(); r.text = line
            r.font.size = Pt(10); r.font.color.rgb = WHITE
        else:
            para(tf2, line, sz=10, c=WHITE, bullet=False)
    # Flèche
    if i < 3:
        txt(sl, "►", cx + Inches(3.05), Inches(2.3), Inches(0.3), Inches(0.4),
            sz=16, c=GRAY, align=PP_ALIGN.CENTER)

# Variables retenues (tableau simplifié)
rect(sl, Inches(0.4), Inches(3.95), Inches(12.5), Inches(2.75), CARD_BG, rnd=True)
rect(sl, Inches(0.4), Inches(3.95), Inches(12.5), Inches(0.4), NAVY, rnd=True)
txt(sl, "Variables Prédictives Retenues",
    Inches(0.6), Inches(3.98), Inches(11.0), Inches(0.35),
    sz=11, bold=True, c=WHITE)

cols = [
    ("Sociodémographiques",
     "Âge, Niveau d'éducation, Région, Milieu de résidence, Richesse, Statut matrimonial, Religion, Emploi"),
    ("Reproductives",
     "Nombre d'enfants nés vivants, Âge à la 1ère naissance, Naissances récentes (5 ans), Grossesses interrompues"),
    ("Accès aux soins",
     "Assurance maladie, Consultation dans les 12 mois, Obstacles aux soins, Électricité au domicile"),
]
col_anims = []
for i, (cat, vars_) in enumerate(cols):
    cx = Inches(0.55 + i * 4.15)
    txt(sl, cat, cx, Inches(4.48), Inches(3.9), Inches(0.35),
        sz=10.5, bold=True, c=BLUE)
    bv = txt(sl, vars_, cx, Inches(4.88), Inches(3.9), Inches(0.72),
             sz=9.5, c=DARK, wrap=True)
    col_anims.append(bv)

foot(sl, pn)
click_fade(sl, col_anims)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 7  –  STATISTIQUES DESCRIPTIVES
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Statistiques Descriptives",
     sub="Profil sociodémographique et reproductif des femmes de l'échantillon",
     badge="PARTIE I", badge_c=BLUE)

# Image prévalence bivariée (gauche)
img_bv = img(sl, "outputs_stat/prevalence_bivariee.png",
             Inches(0.35), Inches(1.55), Inches(7.8), Inches(5.6))

# Stats clés à droite
stats = [
    ("30,2 ans",   "Âge moyen des femmes",            BLUE),
    ("34,1 %",     "Prévalence mortalité infantile",   RED),
    ("3,2",        "Nombre moyen d'enfants nés vivants", GREEN),
    ("21,4 ans",   "Âge médian à la 1ère naissance",  ORANGE),
    ("47 %",       "Niveau secondaire (le plus fréquent)", NAVY),
    ("55 %",       "Milieu rural",                     GRAY),
]
rect(sl, Inches(8.4), Inches(1.55), Inches(4.55), Inches(5.6), CARD_BG, rnd=True)
for i, (val, lab, c) in enumerate(stats):
    cy = Inches(1.72 + i * 0.9)
    rect(sl, Inches(8.55), cy, Inches(4.25), Inches(0.72), WHITE, rnd=True)
    rect(sl, Inches(8.55), cy, Inches(0.06), Inches(0.72), c)
    txt(sl, val, Inches(8.7), cy + Inches(0.06), Inches(1.6), Inches(0.35),
        sz=14, bold=True, c=c)
    txt(sl, lab, Inches(10.35), cy + Inches(0.1), Inches(2.35), Inches(0.52),
        sz=9.5, c=DARK, wrap=True)

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 8  –  ANALYSE BIVARIÉE
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Analyse Bivariée",
     sub="Tests du Khi-deux et t-tests — Association avec la mortalité infantile",
     badge="PARTIE I", badge_c=BLUE)

# Tableau des résultats bivariés
headers = ["Variable", "Test", "p-valeur", "Significativité"]
rows = [
    ["Niveau d'éducation",      "Khi-deux", "< 0,001", "***"],
    ["Richesse du ménage",       "Khi-deux", "< 0,001", "***"],
    ["Région de résidence",      "Khi-deux", "< 0,001", "***"],
    ["Âge à la 1ère naissance",  "Khi-deux", "< 0,001", "***"],
    ["Nombre d'enfants",         "t-test",   "< 0,001", "***"],
    ["Milieu de résidence",      "Khi-deux", "< 0,001", "***"],
    ["Statut matrimonial",       "Khi-deux", "< 0,001", "***"],
    ["Emploi",                   "Khi-deux", "< 0,001", "***"],
    ["Assurance maladie",        "Khi-deux", "< 0,001", "***"],
    ["Grossesses interrompues",  "Khi-deux", "< 0,001", "***"],
    ["Naissances récentes",      "Khi-deux", "< 0,001", "***"],
    ["Religion",                 "Khi-deux", "0,002",   "**"],
    ["Consultation soins",       "Khi-deux", "< 0,001", "***"],
    ["Électricité",              "Khi-deux", "< 0,001", "***"],
]

# En-tête du tableau
col_w = [Inches(3.2), Inches(1.5), Inches(1.5), Inches(1.4)]
col_x = [Inches(0.4), Inches(3.65), Inches(5.2), Inches(6.75)]
row_h = Inches(0.34)

for ci, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
    rect(sl, cx, Inches(1.55), cw, Inches(0.4), NAVY)
    txt(sl, hdr, cx + Inches(0.07), Inches(1.57), cw - Inches(0.1), Inches(0.36),
        sz=10.5, bold=True, c=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    cy = Inches(1.95) + ri * row_h
    bg_r = CARD_BG if ri % 2 == 0 else WHITE
    for ci, (val, cw, cx) in enumerate(zip(row, col_w, col_x)):
        rect(sl, cx, cy, cw, row_h, bg_r,
             line=LGRAY, lw=0.5)
        c_val = DARK
        sz_val = 10
        if ci == 3:
            c_val = GREEN if val == "***" else ORANGE
            sz_val = 11
        txt(sl, val, cx + Inches(0.07), cy + Inches(0.06),
            cw - Inches(0.1), row_h - Inches(0.07),
            sz=sz_val, c=c_val,
            align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

# Légende
txt(sl, "*** p < 0,001   ** p < 0,01   Toutes les 14 variables sont significativement associées à la mortalité infantile.",
    Inches(0.4), Inches(6.88), Inches(8.5), Inches(0.35),
    sz=9.5, italic=True, c=GRAY)

# Résumé clé
rect(sl, Inches(8.45), Inches(1.55), Inches(4.55), Inches(5.6), CARD_BG, rnd=True)
rect(sl, Inches(8.45), Inches(1.55), Inches(4.55), Inches(0.06), GREEN)
txt(sl, "Principaux résultats",
    Inches(8.65), Inches(1.72), Inches(4.1), Inches(0.38),
    sz=12, bold=True, c=NAVY)

bx_res = sl.shapes.add_textbox(Inches(8.65), Inches(2.2), Inches(4.1), Inches(4.8))
tf_r = bx_res.text_frame; tf_r.word_wrap = True
res_items = [
    "14 / 14 variables significatives (p < 0,05)",
    "Variables les plus fortement associées :\n  → Éducation, parité, âge 1ère naissance",
    "Richesse et région : reflet des inégalités socio-économiques",
    "Accès aux soins : fort effet protecteur",
    "Fondement pour la régression multivariée",
]
for j, item in enumerate(res_items):
    if j == 0:
        p = tf_r.paragraphs[0]; r = p.add_run(); r.text = item
        r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = GREEN
    else:
        para(tf_r, item, sz=10.5, c=DARK, bullet=True, symbol='▸')

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 9  –  RÉGRESSION LOGISTIQUE PONDÉRÉE
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Régression Logistique Pondérée",
     sub="Modèle complet intégrant l'ensemble des blocs de variables — n = 10 494 mères",
     badge="PARTIE I", badge_c=BLUE)

# Trois blocs de variables du modèle unique
blocs = [
    ("Bloc 1 : Sociodémographique",
     BLUE,
     "Âge, niveau d'éducation, région, milieu de résidence,\nquintile de richesse, statut matrimonial, religion, emploi"),
    ("Bloc 2 : Reproductif",
     NAVY,
     "Nombre d'enfants nés vivants, âge à la 1ère naissance,\nnaissances récentes (5 ans), grossesses interrompues"),
    ("Bloc 3 : Accès aux soins",
     RGBColor(0x6B, 0x21, 0xA8),
     "Assurance maladie, consultation établissement de soins,\nvisite agent de santé, électricité, obstacles aux soins"),
]
for i, (name, c, vars_) in enumerate(blocs):
    w = Inches(4.15)
    cx = Inches(0.35 + i * 4.32)
    bk = rect(sl, cx, Inches(1.6), w, Inches(2.15), c, rnd=True)
    txt(sl, name, cx + Inches(0.15), Inches(1.7), w - Inches(0.25), Inches(0.6),
        sz=11, bold=True, c=GOLD, wrap=True)
    txt(sl, vars_, cx + Inches(0.15), Inches(2.32), w - Inches(0.25), Inches(1.3),
        sz=9.5, c=WHITE, wrap=True)
    if i < 2:
        txt(sl, "+", cx + w + Inches(0.08), Inches(2.2), Inches(0.3), Inches(0.45),
            sz=22, bold=True, c=GOLD, align=PP_ALIGN.CENTER)

# Paramètres méthodologiques
rect(sl, Inches(0.4), Inches(3.95), Inches(12.5), Inches(2.85), CARD_BG, rnd=True)
rect(sl, Inches(0.4), Inches(3.95), Inches(12.5), Inches(0.4), NAVY, rnd=True)
txt(sl, "Paramètres Méthodologiques",
    Inches(0.6), Inches(3.98), Inches(11.0), Inches(0.35),
    sz=11, bold=True, c=WHITE)

param_cols = [
    ("Pondération",
     ["Poids normalisés : w = V005/10⁶\ndivisé par la moyenne des poids",
      "Écart-type robuste (HC1)",
      "Hypothèse : MCAR sur les poids manquants"]),
    ("Qualité du modèle",
     ["AUC sur données test (25 %)",
      "Pseudo-R² de McFadden et Nagelkerke",
      "Test de Hosmer-Lemeshow (calibration)"]),
    ("Diagnostic",
     ["VIF pour toutes les variables\n(seuil : VIF > 10 = problématique)",
      "Aucune multicolinéarité problématique",
      "Séparation parfaite évitée (sous-groupe mères)"]),
    ("Catégories de référence",
     ["Éducation : Supérieur",
      "Résidence : Urbain",
      "Richesse : Très riche",
      "Âge 1ère naissance : 25 ans et plus"]),
]
for i, (cat, items) in enumerate(param_cols):
    cx = Inches(0.6 + i * 3.15)
    txt(sl, cat, cx, Inches(4.5), Inches(2.95), Inches(0.35),
        sz=10.5, bold=True, c=BLUE)
    bv = sl.shapes.add_textbox(cx, Inches(4.88), Inches(2.95), Inches(1.8))
    tf_v = bv.text_frame; tf_v.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            p = tf_v.paragraphs[0]; r = p.add_run(); r.text = item
            r.font.size = Pt(9.5); r.font.color.rgb = DARK
        else:
            para(tf_v, item, sz=9.5, c=DARK, bullet=True, symbol='▸')

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 10  –  RÉSULTATS – ODDS RATIOS
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Résultats — Facteurs de Risque et Protecteurs",
     sub="Odds ratios ajustés — Modèle 3 complet — n = 10 494 femmes avec enfants",
     badge="PARTIE I", badge_c=BLUE)

# Forest plot (si disponible)
fp = img(sl, "outputs_stat/forest_plot_stat.png",
         Inches(0.35), Inches(1.58), Inches(5.85), Inches(5.6))

# Tableau des OR clés
rect(sl, Inches(6.5), Inches(1.58), Inches(6.5), Inches(5.6), CARD_BG, rnd=True)
rect(sl, Inches(6.5), Inches(1.58), Inches(6.5), Inches(0.42), NAVY, rnd=True)
txt(sl, "Odds Ratios Ajustés — Résultats Clés",
    Inches(6.7), Inches(1.62), Inches(6.1), Inches(0.35),
    sz=10.5, bold=True, c=WHITE)

or_data = [
    ("Facteurs de RISQUE",      None,   None, RED,   True),
    ("Âge 1ère naissance < 18", "5,05", "(4,11 – 6,21)", RED,   False),
    ("Aucun niveau d'éducation", "2,98", "(2,04 – 4,34)", RED,   False),
    ("Âge 1ère naissance 18-19","2,64", "(2,13 – 3,28)", RED,   False),
    ("Éducation primaire",       "2,30", "(1,61 – 3,29)", ORANGE,False),
    ("Éducation secondaire",     "1,49", "(1,05 – 2,11)", ORANGE,False),
    ("Facteurs PROTECTEURS",     None,   None, GREEN, True),
    ("Jamais en union",          "0,69", "(0,55 – 0,85)", GREEN, False),
    ("Emploi actif",             "0,81", "(0,72 – 0,92)", GREEN, False),
    ("Consultation soins (12 m)","0,77", "(0,68 – 0,86)", GREEN, False),
]
row_h2 = Inches(0.46)
for ri, (var, or_v, ci_v, c, is_header) in enumerate(or_data):
    cy = Inches(2.12) + ri * row_h2
    if is_header:
        rect(sl, Inches(6.5), cy, Inches(6.5), row_h2,
             RGBColor(0xD0, 0xD8, 0xDF))
        txt(sl, var, Inches(6.65), cy + Inches(0.08),
            Inches(6.1), row_h2 - Inches(0.1),
            sz=10, bold=True, c=c)
    else:
        bg_r = WHITE if ri % 2 == 0 else CARD_BG
        rect(sl, Inches(6.5), cy, Inches(6.5), row_h2,
             bg_r, line=LGRAY, lw=0.4)
        txt(sl, var, Inches(6.65), cy + Inches(0.1),
            Inches(3.2), row_h2 - Inches(0.1), sz=9.5, c=DARK)
        txt(sl, or_v or "", Inches(9.9), cy + Inches(0.1),
            Inches(0.8), row_h2 - Inches(0.1),
            sz=11, bold=True, c=c, align=PP_ALIGN.CENTER)
        txt(sl, ci_v or "", Inches(10.75), cy + Inches(0.1),
            Inches(2.1), row_h2 - Inches(0.1),
            sz=9.5, italic=True, c=GRAY, align=PP_ALIGN.CENTER)

# AUC en bas droite
rect(sl, Inches(6.5), Inches(6.5), Inches(6.5), Inches(0.65), NAVY, rnd=True)
txt(sl, "AUC du modèle final : 0,87   |   Pseudo-R² Nagelkerke : 0,23   |   Seuil : p < 0,05",
    Inches(6.7), Inches(6.55), Inches(6.1), Inches(0.55),
    sz=10, bold=True, c=GOLD, align=PP_ALIGN.CENTER)

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 11  –  PERFORMANCE DU MODÈLE STATISTIQUE
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Performance du Modèle Statistique",
     sub="Courbe ROC et indicateurs de qualité du modèle logistique",
     badge="PARTIE I", badge_c=BLUE)

# Courbe ROC
rc = img(sl, "outputs_stat/roc_curve_stat.png",
         Inches(0.4), Inches(1.55), Inches(6.0), Inches(5.6))

# Métriques de performance
rect(sl, Inches(6.75), Inches(1.55), Inches(6.2), Inches(5.6), CARD_BG, rnd=True)
rect(sl, Inches(6.75), Inches(1.55), Inches(6.2), Inches(0.45), BLUE, rnd=True)
txt(sl, "Indicateurs de Qualité du Modèle",
    Inches(6.95), Inches(1.58), Inches(5.8), Inches(0.38),
    sz=11, bold=True, c=WHITE)

metrics2 = [
    ("AUC",                 "0,87",   "Discrimination excellente",   BLUE),
    ("Pseudo-R² McFadden",  "0,14",   "Contribution relative modérée",NAVY),
    ("Pseudo-R² Nagelkerke","0,23",   "Variance expliquée estimée",  BLUE),
    ("Log-vraisemblance",   "Model 3 meilleur\n(AIC le plus faible)",
                            "Comparaison modèles emboités",        GREEN),
    ("VIF max",             "< 5",    "Pas de multicolinéarité",     GREEN),
    ("Effectif",            "10 494", "Mères (au moins 1 naissance vivante)", GRAY),
]
for i, (name, val, desc, c) in enumerate(metrics2):
    cy = Inches(2.18) + i * Inches(0.78)
    rect(sl, Inches(6.9), cy, Inches(5.9), Inches(0.66), WHITE, line=LGRAY, lw=0.5)
    rect(sl, Inches(6.9), cy, Inches(0.06), Inches(0.66), c)
    txt(sl, name, Inches(7.05), cy + Inches(0.06), Inches(2.1), Inches(0.3),
        sz=10, bold=True, c=NAVY)
    txt(sl, val, Inches(9.25), cy + Inches(0.04), Inches(1.5), Inches(0.38),
        sz=13, bold=True, c=c, align=PP_ALIGN.CENTER)
    txt(sl, desc, Inches(10.85), cy + Inches(0.1), Inches(1.8), Inches(0.48),
        sz=8.5, italic=True, c=GRAY, wrap=True)

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 12  –  SÉPARATEUR PARTIE II
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide(bg=NAVY2)

rect(sl, 0, 0, SW, SH, NAVY2)
rect(sl, 0, Inches(3.1), SW, Inches(1.3), BLUE)
rect(sl, 0, Inches(3.1), SW, Inches(0.06), GOLD)
rect(sl, 0, Inches(4.34), SW, Inches(0.06), GOLD)

txt(sl, "II", Inches(0.5), Inches(0.5), Inches(2.5), Inches(3.5),
    sz=180, bold=True, c=RGBColor(0x2C, 0x5F, 0x8A))

b_p = txt(sl, "PARTIE II",
          Inches(3.0), Inches(3.15), Inches(9.0), Inches(0.55),
          sz=22, bold=True, c=GOLD)

b_t = txt(sl, "Machine Learning",
          Inches(3.0), Inches(3.62), Inches(9.0), Inches(0.72),
          sz=38, bold=True, c=WHITE)

txt(sl,
    "Pipeline ML  •  Six algorithmes  •  Courbes d'apprentissage  •  Comparaison  •  SHAP  •  Application",
    Inches(3.0), Inches(4.45), Inches(9.8), Inches(0.42),
    sz=12, italic=True, c=RGBColor(0xAE, 0xC6, 0xE0))

foot(sl, pn)
auto_fade(sl, [b_p, b_t], start_delay=200, step=400)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 13  –  PIPELINE DE PRÉTRAITEMENT
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Pipeline de Prétraitement Machine Learning",
     sub="Traitement différencié selon la nature des variables — Seed fixé à 42",
     badge="PARTIE II", badge_c=NAVY)

# Trois colonnes de variables
var_groups = [
    ("Variables Numériques\n(6 variables)",
     SKY,
     "Âge, Nombre d'enfants, Âge 1ère naissance,\nNaissances récentes, Grossesses interrompues,\nNombre de membres du ménage",
     "Imputation → Médiane\nMise à l'échelle → StandardScaler"),
    ("Variables Ordinales\n(14 variables)",
     BLUE,
     "Éducation, Richesse, Région, Religion,\nMilieu, Statut matrimonial, Emploi,\nAssurance, Obstacles aux soins…",
     "Imputation → Mode\n(catégorie la plus fréquente)"),
    ("Variables Catégorielles\n(2 variables)",
     NAVY,
     "Électricité, Consultation dans les 12 mois",
     "Imputation → Mode\nEncodage → OneHotEncoder\nhandle unknown = ignore"),
]
for i, (title, c, vars_, steps_) in enumerate(var_groups):
    cx = Inches(0.35 + i * 4.32)
    rect(sl, cx, Inches(1.6), Inches(4.1), Inches(3.6), c, rnd=True)
    txt(sl, title, cx + Inches(0.15), Inches(1.72), Inches(3.85), Inches(0.55),
        sz=11, bold=True, c=GOLD, wrap=True)
    txt(sl, vars_, cx + Inches(0.15), Inches(2.32), Inches(3.85), Inches(1.5),
        sz=9.5, c=WHITE, wrap=True)
    rect(sl, cx + Inches(0.1), Inches(3.82), Inches(3.9), Inches(0.06),
         GOLD)
    txt(sl, steps_, cx + Inches(0.15), Inches(3.95), Inches(3.85), Inches(1.1),
        sz=9.5, italic=True, c=RGBColor(0xAE, 0xC6, 0xE0), wrap=True)

# Split et validation
rect(sl, Inches(0.4), Inches(5.45), Inches(12.5), Inches(1.65), CARD_BG, rnd=True)
rect(sl, Inches(0.4), Inches(5.45), Inches(12.5), Inches(0.4), BLUE, rnd=True)
txt(sl, "Stratégie de Validation",
    Inches(0.6), Inches(5.48), Inches(11.0), Inches(0.35),
    sz=11, bold=True, c=WHITE)

split_items = [
    ("75 % Entraînement\n(11 007 obs.)", NAVY),
    ("25 % Test\n(3 670 obs.)", BLUE),
    ("5 plis stratifiés\n(Validation croisée)", SKY),
    ("Classe déséquilibrée\nclass weight = balanced", ORANGE),
    ("Seuil optimal\nmaximisant le F1-Score", GREEN),
]
for i, (text_, c) in enumerate(split_items):
    cx = Inches(0.6 + i * 2.5)
    rect(sl, cx, Inches(5.98), Inches(2.35), Inches(0.88), c, rnd=True)
    txt(sl, text_, cx + Inches(0.12), Inches(6.02), Inches(2.15), Inches(0.8),
        sz=9.5, bold=True, c=WHITE, align=PP_ALIGN.CENTER, wrap=True)

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 14  –  CINQ MODÈLES DE MACHINE LEARNING
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Six Algorithmes de Machine Learning",
     sub="Comparaison d'approches linéaires, d'ensemble et de gradient boosting avancé",
     badge="PARTIE II", badge_c=NAVY)

models_ml = [
    ("Régression\nLogistique",
     BLUE, "C = 0,1 | L2\nmax iter = 500",
     "Référence linéaire\nInterprétable"),
    ("Random\nForest",
     GREEN, "200 arbres\nmax depth = 10\nmin samples = 10",
     "Ensemble bagging\nVariance réduite"),
    ("Gradient\nBoosting",
     ORANGE, "200 estimateurs\nlr = 0,05\nsubsample = 0,8",
     "Boosting séquentiel\nSklearn natif"),
    ("XGBoost",
     RED, "300 estimateurs\nmax depth = 5\nlr = 0,05",
     "Optimisé CPU/GPU\nMeilleur modèle"),
    ("LightGBM",
     NAVY, "300 estimateurs\ncroissance feuilles\nlr = 0,05",
     "Très rapide\nGrandes données"),
    ("CatBoost",
     RGBColor(0xFF, 0xC0, 0x07), "300 iterations\ndepth = 6\nauto_class_weights",
     "Gère nativement\nle déséquilibre"),
]
anim_cards = []
card_w = Inches(2.05)
for i, (name, c, params, desc) in enumerate(models_ml):
    cx = Inches(0.25 + i * 2.18)
    bk = rect(sl, cx, Inches(1.6), card_w, Inches(5.5), CARD_BG, rnd=True)
    rect(sl, cx, Inches(1.6), card_w, Inches(0.06), c)
    anim_cards.append(bk)
    txt(sl, name, cx + Inches(0.12), Inches(1.72), card_w - Inches(0.2), Inches(0.65),
        sz=11.5, bold=True, c=c, wrap=True)
    rect(sl, cx + card_w - Inches(0.62), Inches(1.68), Inches(0.52), Inches(0.52), c, rnd=True)
    txt(sl, str(i+1), cx + card_w - Inches(0.62), Inches(1.7), Inches(0.52), Inches(0.48),
        sz=13, bold=True, c=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "Hyperparamètres :",
        cx + Inches(0.12), Inches(2.45), card_w - Inches(0.2), Inches(0.28),
        sz=8, bold=True, c=NAVY)
    txt(sl, params, cx + Inches(0.12), Inches(2.75), card_w - Inches(0.2), Inches(0.88),
        sz=8.5, c=DARK, wrap=True)
    rect(sl, cx + Inches(0.12), Inches(3.7), card_w - Inches(0.2), Inches(0.03), LGRAY)
    txt(sl, desc, cx + Inches(0.12), Inches(3.78), card_w - Inches(0.2), Inches(0.65),
        sz=9, italic=True, c=GRAY, wrap=True)
    rect(sl, cx, Inches(4.65), card_w, Inches(2.12), c)
    txt(sl, "Déséquilibre\nY=0 : 75 %  |  Y=1 : 25 %\nclass weight = balanced",
        cx + Inches(0.08), Inches(4.72), card_w - Inches(0.12), Inches(2.0),
        sz=8, c=WHITE, wrap=True)

foot(sl, pn)
click_fade(sl, anim_cards)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 15  –  COURBES D'APPRENTISSAGE
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Courbes d'Apprentissage",
     sub="AUC entraînement vs validation croisée (5-fold) en fonction de la taille de l'échantillon",
     badge="PARTIE II", badge_c=NAVY)

# Image : grille des 6 courbes d'apprentissage
lc_grid = img(sl, "outputs_ml/learning_curves_all_models.png",
              Inches(0.35), Inches(1.55), Inches(9.0), Inches(5.5))

# Interprétation à droite
rect(sl, Inches(9.65), Inches(1.55), Inches(3.55), Inches(5.5), CARD_BG, rnd=True)
rect(sl, Inches(9.65), Inches(1.55), Inches(3.55), Inches(0.42), NAVY, rnd=True)
txt(sl, "Lecture des courbes",
    Inches(9.85), Inches(1.58), Inches(3.2), Inches(0.36),
    sz=10.5, bold=True, c=WHITE)

bx_interp = sl.shapes.add_textbox(Inches(9.82), Inches(2.1), Inches(3.2), Inches(4.75))
tf_i = bx_interp.text_frame; tf_i.word_wrap = True
interp_items = [
    "Courbe bleue : score sur données d'entraînement",
    "Courbe rouge : score sur validation croisée",
    "Un faible écart final indique une bonne généralisation",
    "Convergence des deux courbes = pas de surapprentissage",
    "Tous les modèles convergent > 0,80 AUC",
    "XGBoost : gap < 0,05 — excellent équilibre biais / variance",
    "Le modèle bénéficie de données supplémentaires (courbe val. croissante)",
]
for j, item in enumerate(interp_items):
    if j == 0:
        p = tf_i.paragraphs[0]; r = p.add_run(); r.text = item
        r.font.size = Pt(9.5); r.font.color.rgb = DARK
    else:
        para(tf_i, item, sz=9.5, c=DARK, bullet=True, symbol='▸')

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 16  –  COMPARAISON DES PERFORMANCES
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Comparaison des Performances",
     sub="Évaluation sur l'ensemble de test (25 %) — 3 670 observations",
     badge="PARTIE II", badge_c=NAVY)

# Image comparaison ROC
rc_ml = img(sl, "outputs_ml/roc_comparison_ml.png",
            Inches(0.4), Inches(1.55), Inches(5.8), Inches(5.6))

# Tableau des performances
rect(sl, Inches(6.55), Inches(1.55), Inches(6.4), Inches(5.6), CARD_BG, rnd=True)
rect(sl, Inches(6.55), Inches(1.55), Inches(6.4), Inches(0.42), NAVY, rnd=True)
txt(sl, "Métriques sur l'Ensemble de Test",
    Inches(6.75), Inches(1.58), Inches(6.1), Inches(0.35),
    sz=10.5, bold=True, c=WHITE)

perf_headers = ["Modèle", "AUC CV", "AUC Test", "F1", "Rappel"]
perf_data = [
    ("Régression Log.", "0,877", "0,874", "0,602", "0,738", False),
    ("Random Forest",   "0,879", "0,874", "0,609", "0,711", False),
    ("Gradient Boosting","0,887","0,887", "0,635", "0,695", False),
    ("XGBoost",         "0,885", "0,889", "0,643", "0,720", True),
    ("LightGBM",        "0,884", "0,889", "0,637", "0,748", False),
    ("CatBoost",        "0,885", "0,887", "0,631", "0,720", False),
]
ph_w = [Inches(1.75), Inches(0.85), Inches(0.9), Inches(0.8), Inches(0.85)]
ph_x = [Inches(6.62), Inches(8.42), Inches(9.32), Inches(10.27), Inches(11.12)]
row_h_perf = Inches(0.62)

for ci, (h, cw, cx) in enumerate(zip(perf_headers, ph_w, ph_x)):
    rect(sl, cx, Inches(2.08), cw, Inches(0.38), NAVY)
    txt(sl, h, cx + Inches(0.04), Inches(2.1), cw - Inches(0.07), Inches(0.34),
        sz=9, bold=True, c=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(perf_data):
    *vals, is_best = row
    cy = Inches(2.46) + ri * row_h_perf
    bg_r = RGBColor(0xD4, 0xED, 0xDA) if is_best else (CARD_BG if ri % 2 == 0 else WHITE)
    for ci, (val, cw, cx) in enumerate(zip(vals, ph_w, ph_x)):
        rect(sl, cx, cy, cw, row_h_perf, bg_r, line=LGRAY, lw=0.4)
        c_v = DARK
        if ci > 0:
            c_v = RED if is_best else DARK
        bld = (ci == 2 or ci == 3) and is_best
        txt(sl, val, cx + Inches(0.04), cy + Inches(0.14),
            cw - Inches(0.07), Inches(0.34),
            sz=10 if ci > 0 else 9,
            bold=bld, c=c_v, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
    if is_best:
        txt(sl, "★", ph_x[-1] + ph_w[-1] + Inches(0.04), cy + Inches(0.18),
            Inches(0.45), Inches(0.3), sz=12, bold=True, c=GREEN)

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 16  –  MODÈLE CHAMPION : XGBOOST
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Modèle Champion — XGBoost",
     sub="Importance des variables et interprétabilité via SHAP",
     badge="PARTIE II", badge_c=NAVY)

# Image SHAP
sh = img(sl, "outputs_ml/shap_importance_XGBoost.png",
         Inches(0.35), Inches(1.55), Inches(6.2), Inches(5.6))

# Image importance features
fi = img(sl, "outputs_ml/feature_importance_XGBoost.png",
         Inches(6.8), Inches(1.55), Inches(3.0), Inches(2.7))

# Métriques XGBoost
metrics_xgb = [
    ("AUC",     "0,889", GREEN),
    ("F1",      "0,642", BLUE),
    ("Rappel",  "0,720", NAVY),
    ("Précision","0,580", ORANGE),
]
for i, (lbl, val, c) in enumerate(metrics_xgb):
    cx = Inches(6.8 + i * 1.56)
    rect(sl, cx, Inches(4.52), Inches(1.45), Inches(0.9), c, rnd=True)
    txt(sl, lbl, cx + Inches(0.07), Inches(4.56), Inches(1.35), Inches(0.28),
        sz=8.5, bold=True, c=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, val, cx + Inches(0.07), Inches(4.82), Inches(1.35), Inches(0.48),
        sz=18, bold=True, c=GOLD, align=PP_ALIGN.CENTER)

# Top features
rect(sl, Inches(6.8), Inches(5.58), Inches(6.1), Inches(1.55), CARD_BG, rnd=True)
rect(sl, Inches(6.8), Inches(5.58), Inches(6.1), Inches(0.38), NAVY, rnd=True)
txt(sl, "Top Variables Prédictives",
    Inches(7.0), Inches(5.61), Inches(5.8), Inches(0.32),
    sz=10, bold=True, c=WHITE)

top_feats = [
    ("Nombre d'enfants nés vivants", "23,1 %", RED),
    ("Âge de la femme",              " 4,2 %", ORANGE),
    ("Niveau d'éducation",           " 4,1 %", BLUE),
    ("Naissances récentes (5 ans)",  " 2,8 %", NAVY),
    ("Âge à la 1ère naissance",      " 2,7 %", GREEN),
]
for i, (feat, pct, c) in enumerate(top_feats):
    cx2 = Inches(6.95)
    cy2 = Inches(6.05) + i * Inches(0.2)
    txt(sl, feat, cx2, cy2, Inches(3.8), Inches(0.2),
        sz=8.5, c=DARK)
    # Barre de progression
    bar_w = float(pct.strip().replace('%', '').replace(',', '.').strip()) / 25.0 * Inches(1.9)
    rect(sl, Inches(10.85), cy2 + Inches(0.02), Inches(1.9), Inches(0.16), LGRAY)
    rect(sl, Inches(10.85), cy2 + Inches(0.02), bar_w, Inches(0.16), c)
    txt(sl, pct.strip(), Inches(12.82), cy2, Inches(0.6), Inches(0.2),
        sz=8, bold=True, c=c, align=PP_ALIGN.RIGHT)

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 17  –  APPLICATION STREAMLIT
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Application de Prédiction — Streamlit",
     sub="Interface interactive d'estimation du risque de mortalité infantile",
     badge="PARTIE II", badge_c=NAVY)

# Schéma de l'interface
rect(sl, Inches(0.4), Inches(1.6), Inches(12.5), Inches(5.5), CARD_BG, rnd=True)

# Sidebar simulée
rect(sl, Inches(0.55), Inches(1.75), Inches(3.1), Inches(5.2), NAVY, rnd=True)
txt(sl, "Saisie des\nInformations Patient",
    Inches(0.7), Inches(1.85), Inches(2.85), Inches(0.6),
    sz=10, bold=True, c=GOLD, wrap=True)

sidebar_items = [
    ("Profil Sociodémographique",
     "Âge, éducation, région, richesse, religion, emploi"),
    ("Historique Reproductif",
     "Nombre d'enfants, âge 1ère naissance, naissances récentes"),
    ("Accès aux Soins",
     "Assurance maladie, consultations, obstacles aux soins"),
]
for i, (cat, det) in enumerate(sidebar_items):
    cy = Inches(2.58) + i * Inches(1.38)
    rect(sl, Inches(0.65), cy, Inches(2.85), Inches(1.15),
         RGBColor(0x2A, 0x58, 0x8A), rnd=True)
    txt(sl, cat, Inches(0.75), cy + Inches(0.08), Inches(2.65), Inches(0.38),
        sz=9, bold=True, c=GOLD)
    txt(sl, det, Inches(0.75), cy + Inches(0.48), Inches(2.65), Inches(0.62),
        sz=8.5, c=WHITE, wrap=True)

txt(sl, "► Estimer le Risque",
    Inches(0.65), Inches(6.45), Inches(2.9), Inches(0.28),
    sz=10, bold=True, c=WHITE, align=PP_ALIGN.CENTER)
rect(sl, Inches(0.65), Inches(6.44), Inches(2.9), Inches(0.3), SKY, rnd=True)
txt(sl, "► Estimer le Risque",
    Inches(0.65), Inches(6.45), Inches(2.9), Inches(0.28),
    sz=10, bold=True, c=WHITE, align=PP_ALIGN.CENTER)

# Résultats simulés – Machine Learning uniquement
rect(sl, Inches(3.85), Inches(1.75), Inches(8.8), Inches(5.2), WHITE, line=LGRAY, lw=0.8)
txt(sl, "Machine Learning — XGBoost (AUC = 0,889)",
    Inches(4.0), Inches(1.82), Inches(8.5), Inches(0.38),
    sz=11.5, bold=True, c=NAVY, align=PP_ALIGN.CENTER)

# Jauge de risque simulée
rect(sl, Inches(4.5), Inches(2.35), Inches(7.5), Inches(0.22),
     RGBColor(0xDE, 0xE2, 0xE6), rnd=True)
rect(sl, Inches(4.5), Inches(2.35), Inches(4.35), Inches(0.22),
     RED, rnd=True)
txt(sl, "58,7 %", Inches(8.7), Inches(2.2), Inches(1.2), Inches(0.32),
    sz=11, bold=True, c=RED)

rect(sl, Inches(5.5), Inches(2.72), Inches(5.6), Inches(0.95), CARD_R, rnd=True)
txt(sl, "Probabilité : 58,7 %",
    Inches(5.5), Inches(2.88), Inches(5.6), Inches(0.38),
    sz=12, bold=True, c=RED, align=PP_ALIGN.CENTER)
txt(sl, "RISQUE ÉLEVÉ",
    Inches(4.8), Inches(3.82), Inches(7.0), Inches(0.52),
    sz=16, bold=True, c=RED, align=PP_ALIGN.CENTER)
txt(sl,
    "Suivi médical urgent requis  ·  Éducation à la santé maternelle  ·  Planification familiale conseillée",
    Inches(4.2), Inches(4.55), Inches(8.1), Inches(0.5),
    sz=9.5, c=DARK, wrap=True, align=PP_ALIGN.CENTER)

# Métriques
metrics_sim = [("AUC", "0,889", GREEN), ("F1", "0,643", BLUE), ("Rappel", "0,720", NAVY)]
for i, (lbl, v, c_) in enumerate(metrics_sim):
    cx_m = Inches(5.0 + i * 2.3)
    rect(sl, cx_m, Inches(5.2), Inches(2.0), Inches(0.65), c_, rnd=True)
    txt(sl, lbl, cx_m + Inches(0.08), Inches(5.24), Inches(1.85), Inches(0.25),
        sz=8.5, bold=True, c=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, v, cx_m + Inches(0.08), Inches(5.46), Inches(1.85), Inches(0.35),
        sz=14, bold=True, c=GOLD, align=PP_ALIGN.CENTER)

# Recommandations en bas
rect(sl, Inches(3.85), Inches(6.25), Inches(8.8), Inches(0.55), CARD_BG)
txt(sl,
    "Résultats confirmés par les valeurs SHAP — interprétabilité garantie",
    Inches(3.95), Inches(6.32), Inches(8.6), Inches(0.42),
    sz=9.5, italic=True, c=GRAY, align=PP_ALIGN.CENTER)

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE 18  –  SYNTHÈSE ET CONVERGENCE
# ═════════════════════════════════════════════════════════════════════════════
sl, pn = new_slide()

head(sl, "Synthèse — Convergence des Deux Approches",
     sub="Résultats robustes et cohérents entre la statistique et le Machine Learning")

rect(sl, Inches(0.4), Inches(1.6), Inches(12.5), Inches(2.2), CARD_BG, rnd=True)
rect(sl, Inches(0.4), Inches(1.6), Inches(12.5), Inches(0.42), BLUE, rnd=True)
txt(sl, "Facteurs de Risque Identifiés par les Deux Approches",
    Inches(0.6), Inches(1.63), Inches(11.5), Inches(0.36),
    sz=11, bold=True, c=WHITE)

convergence = [
    ("Âge précoce à la 1ère naissance\n(< 18 ans)",
     "OR = 5,05 ***", "Importance : 2,7 %", RED),
    ("Faible niveau d'éducation\n(aucun ou primaire)",
     "OR = 2,98 ***", "Importance : 4,1 %", ORANGE),
    ("Parité élevée\n(nombre d'enfants)",
     "Modèle 3 sig.", "Importance : 23,1 %", BLUE),
    ("Accès limité aux soins\n(consultation, assurance)",
     "OR = 0,77 ***\n(effet protecteur)", "Importance élevée", GREEN),
]
for i, (factor, stat_r, ml_r, c) in enumerate(convergence):
    cx = Inches(0.6 + i * 3.1)
    rect(sl, cx, Inches(2.15), Inches(2.9), Inches(1.5), WHITE, line=c, lw=1.5, rnd=True)
    txt(sl, factor, cx + Inches(0.12), Inches(2.22), Inches(2.68), Inches(0.55),
        sz=9.5, bold=True, c=c, wrap=True)
    txt(sl, stat_r, cx + Inches(0.12), Inches(2.85), Inches(2.68), Inches(0.3),
        sz=8.5, italic=True, c=NAVY, wrap=True)
    txt(sl, ml_r, cx + Inches(0.12), Inches(3.15), Inches(2.68), Inches(0.3),
        sz=8.5, italic=True, c=GRAY, wrap=True)

# Performance comparative
rect(sl, Inches(0.4), Inches(3.95), Inches(5.9), Inches(2.8), CARD_BG, rnd=True)
rect(sl, Inches(0.4), Inches(3.95), Inches(5.9), Inches(0.4), BLUE, rnd=True)
txt(sl, "Performance Comparative",
    Inches(0.6), Inches(3.98), Inches(5.6), Inches(0.35),
    sz=10.5, bold=True, c=WHITE)
comp_data = [
    ("Logist. Régression (baseline)", "AUC = 0,874", BLUE),
    ("Gradient Boosting",             "AUC = 0,887", ORANGE),
    ("XGBoost / LightGBM (meilleurs)","AUC = 0,889", GREEN),
]
for i, (name, val, c) in enumerate(comp_data):
    cy = Inches(4.5) + i * Inches(0.65)
    rect(sl, Inches(0.55), cy, Inches(5.6), Inches(0.56), WHITE, line=LGRAY, lw=0.4)
    txt(sl, name, Inches(0.68), cy + Inches(0.1), Inches(3.0), Inches(0.36),
        sz=10, c=DARK)
    txt(sl, val, Inches(3.85), cy + Inches(0.1), Inches(2.1), Inches(0.36),
        sz=11, bold=True, c=c, align=PP_ALIGN.CENTER)

# Recommandations
rect(sl, Inches(6.55), Inches(3.95), Inches(6.4), Inches(2.8), CARD_BG, rnd=True)
rect(sl, Inches(6.55), Inches(3.95), Inches(6.4), Inches(0.4), NAVY, rnd=True)
txt(sl, "Recommandations de Santé Publique",
    Inches(6.75), Inches(3.98), Inches(6.1), Inches(0.35),
    sz=10.5, bold=True, c=WHITE)
bx_recs = sl.shapes.add_textbox(Inches(6.72), Inches(4.48), Inches(6.1), Inches(2.2))
tf_recs  = bx_recs.text_frame; tf_recs.word_wrap = True
recs = [
    "Investir dans l'éducation des filles (réduction du risque de 50 à 70 %)",
    "Renforcer l'accès aux consultations prénatales, surtout en zone rurale",
    "Promouvoir la planification familiale et retarder la première grossesse",
    "Cibler les régions et ménages défavorisés (approche différenciée)",
]
for j, rec in enumerate(recs):
    if j == 0:
        p = tf_recs.paragraphs[0]; r = p.add_run(); r.text = rec
        r.font.size = Pt(10); r.font.color.rgb = DARK
    else:
        para(tf_recs, rec, sz=10, c=DARK, bullet=True, symbol='▸')

foot(sl, pn)


# ═════════════════════════════════════════════════════════════════════════════
#  DIAPOSITIVE (BONUS) 18 → devient 19 si on compte, mais on le garde à 18
#  On écrase le compteur pour garder N = 18
# ═════════════════════════════════════════════════════════════════════════════
# Diapositive 18 = Conclusion finale (la 18ème slide réelle)

# On a déjà créé 18 slides. La dernière étant la synthèse.
# On rajoute la conclusion comme 18ème (on colle synthèse et conclusion)
# → On ignore cette notice et on vérifie le compteur

assert _pg[0] == 19, f"Nombre de slides : {_pg[0]} (attendu 19)"

# ═════════════════════════════════════════════════════════════════════════════
#  SAUVEGARDE
# ═════════════════════════════════════════════════════════════════════════════
out = os.path.join(BASE_DIR, "Presentation_Mortalite_Infantile_Cameroun.pptx")
prs.save(out)
print(f"Présentation sauvegardée : {out}")
print(f"Nombre de diapositives   : {len(prs.slides)}")
